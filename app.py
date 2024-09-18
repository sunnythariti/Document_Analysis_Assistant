import openpyxl
import sys
from pptx import Presentation
import pdfplumber
from docx import Document
from flask import Flask, render_template, request, jsonify, session, redirect, url_for
from werkzeug.utils import secure_filename
from PIL import Image
import os
import base64
from openai import OpenAI
import io
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r'/opt/homebrew/bin/tesseract'

# Increase recursion limit
sys.setrecursionlimit(2000)

# Initialize OpenAI Client
client = OpenAI(api_key='Put_Your_API_Here')

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(__file__), 'uploads')
app.config['MAX_CONTENT_LENGTH'] = 256 * 1024 * 1024
app.secret_key = os.environ.get('FLASK_SECRET_KEY', 'default_secret_key')

ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'xlsx', 'pptx'}

# Global variable to store cached file content
cached_file_content = None

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def estimate_tokens(text):
    """Rough estimate of tokens in a string."""
    return len(text.split())

@app.route('/')
def index():
    if 'logged_in' in session and session['logged_in']:
        return render_template('index.html')
    return redirect(url_for('login'))

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        if request.form['username'] == 'up_to_you' and request.form['password'] == 'put_password_here':
            session['logged_in'] = True
            return redirect(url_for('index'))
        else:
            return render_template('login.html', error='Invalid username or password')
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))

@app.route('/reset', methods=['POST'])
def reset_file():
    if 'logged_in' not in session or not session['logged_in']:
        return jsonify({'error': 'Not logged in'})

    # Clear all file-related data from the session
    session.pop('file_content', None)
    session.pop('filename', None)

    # Delete all files in the upload folder
    for filename in os.listdir(app.config['UPLOAD_FOLDER']):
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        try:
            if os.path.isfile(file_path):
                os.unlink(file_path)
        except Exception as e:
            print(f"Error deleting {file_path}: {e}")

    # Clear any cached data in memory
    global cached_file_content
    cached_file_content = None

    return jsonify({'success': 'File reset successfully'})

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'logged_in' not in session or not session['logged_in']:
        return jsonify({'error': 'Not logged in'})
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'})
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'})
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        try:
            file_content = read_file(file_path)
            session['file_content'] = file_content
            session['filename'] = filename  # Store the filename in the session
            global cached_file_content
            cached_file_content = file_content
            return jsonify({'success': 'File uploaded and read successfully'})
        except Exception as e:
            return jsonify({'error': f'Error reading file: {str(e)}'})
    return jsonify({'error': 'File type not allowed'})

def read_excel_file(file_path):
    workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    all_data = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        sheet_data = [f"Sheet: {sheet_name}"]

        for row in sheet.iter_rows(values_only=True):
            row_data = [str(cell) if cell is not None else "" for cell in row]
            sheet_data.append(", ".join(row_data))

        all_data.extend(sheet_data)
        all_data.append("")  # Add a blank line between sheets

    return "\n".join(all_data)

def read_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    try:
        if file_extension == '.txt':
            with open(file_path, 'r') as file:
                return file.read()
        elif file_extension == '.pdf':
            with pdfplumber.open(file_path) as pdf:
                return "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
        elif file_extension == '.docx':
            doc = Document(file_path)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        elif file_extension == '.xlsx':
            return read_excel_file(file_path)
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
        elif file_extension in ['.png', '.jpg', '.jpeg']:
            with Image.open(file_path) as img:
                # Extract text using OCR
                text = pytesseract.image_to_string(img)

                # Get image metadata
                metadata = f"Image file: {os.path.basename(file_path)}\nFormat: {img.format}\nSize: {img.size}\nMode: {img.mode}"

                if text.strip():
                    return f"{metadata}\n\nExtracted text from image:\n{text}"
                else:
                    return f"{metadata}\n\nNo text could be extracted from this image."
        else:
            return "Unsupported file type"
    except Exception as e:
        print(f"Error reading file: {str(e)}")
        return f"Error reading file: {str(e)}"

@app.route('/chat', methods=['POST'])
def chat():
    if 'logged_in' not in session or not session['logged_in']:
        return jsonify({'error': 'Not logged in'})

    user_message = request.json.get('message')

    # Always read the file content from the file system or use cached content
    global cached_file_content
    file_content = cached_file_content

    if not file_content and 'filename' in session:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], session['filename'])
        if os.path.exists(file_path):
            file_content = read_file(file_path)
            cached_file_content = file_content

    if not file_content:
        return jsonify({'error': 'No file content available. Please upload a file.'})

    if not user_message:
        return jsonify({'error': 'No message provided'})

    try:
        system_message = "You are a helpful assistant. The following is the content of a file. Please use this information to answer the user's questions."

        # Estimate token counts
        system_tokens = estimate_tokens(system_message)
        user_message_tokens = estimate_tokens(user_message)
        file_content_tokens = estimate_tokens(file_content)

        # Reserve tokens for the response
        max_response_tokens = 2000

        # Calculate total tokens
        max_total_tokens = 32768
        total_tokens = system_tokens + user_message_tokens + file_content_tokens + max_response_tokens

        # Truncate file content if necessary
        if total_tokens > max_total_tokens:
            tokens_to_remove = total_tokens - max_total_tokens
            truncation_ratio = 1 - (tokens_to_remove / file_content_tokens)
            truncated_content = file_content[:int(len(file_content) * truncation_ratio)]
        else:
            truncated_content = file_content

        user_content = f"File content:\n{truncated_content}\n\nUser question: {user_message}"

        response = client.chat.completions.create(
            model="gpt-4o",  # Optional to select the model
            messages=[
                {"role": "system", "content": system_message},
                {"role": "user", "content": user_content}
            ],
            max_tokens=max_response_tokens
        )
        return jsonify({'response': response.choices[0].message.content})
    except Exception as e:
        print(f"Error in chat: {str(e)}")
        return jsonify({'error': f"An error occurred: {str(e)}"})

@app.route('/view_file_content')
def view_file_content():
    global cached_file_content
    if cached_file_content:
        return jsonify({'file_content': cached_file_content})
    return jsonify({'error': 'No file content available'})

@app.route('/debug', methods=['GET'])
def debug():
    return jsonify({
        'session': dict(session),
        'upload_folder': app.config['UPLOAD_FOLDER'],
        'allowed_extensions': ALLOWED_EXTENSIONS,
        'has_cached_content': cached_file_content is not None
    })

if __name__ == '__main__':
    if not os.path.exists(app.config['UPLOAD_FOLDER']):
        os.makedirs(app.config['UPLOAD_FOLDER'])
    app.run(debug=True)